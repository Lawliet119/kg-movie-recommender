from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SRC = Path(r"F:\Download\KGenSam cho CRS.pptx")
OUT = Path(r"F:\Download\KGenSam cho CRS - bo sung paper vs code.pptx")
MD_OUT = Path(r"F:\Download\KGenSam paper vs code - bang so sanh.md")


NAVY = RGBColor(20, 33, 61)
BLUE = RGBColor(0, 142, 204)
TEAL = RGBColor(18, 166, 145)
ORANGE = RGBColor(234, 145, 37)
GRAY = RGBColor(95, 99, 104)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)


COMPARISON_ROWS = [
    (
        "Knowledge Graph",
        "KB4REC + Freebase external KG, heterogeneous graph item-entity.",
        "MovieLens + TMDB external KG. Relations: genre, year, director, cast, tag, language.",
        "Implemented demo-scale",
    ),
    (
        "User Feedback",
        "Binary feedback cho item/attribute trong multi-turn CRS.",
        "Accept/Reject attribute; Looks good/Not for me item feedback.",
        "Implemented",
    ),
    (
        "Interact Policy Network",
        "DQN-style policy quyet dinh ASK hay RECOMMEND moi turn.",
        "rollout_dqn train bang MovieLens-derived simulator; co min 3 asks de demo workflow ro.",
        "Implemented demo RL",
    ),
    (
        "Active Sampler",
        "RL sampler chon attribute/active graph tot nhat de hoi.",
        "rollout_gcn + entropy/centrality/FM uncertainty de chon cau hoi.",
        "Implemented demo RL",
    ),
    (
        "Negative Sampler",
        "KG-enhanced negative sampling tao negative samples chat luong cao cho recommender.",
        "learned_linear_policy sampler hoc tu MovieLens feedback + KG similarity features.",
        "Implemented demo policy",
    ),
    (
        "Recommender",
        "FM/BPR recommender duoc update voi sampled negatives va KG signals.",
        "FM/BPR + KG propagation + KG embedding similarity, hybrid ranking.",
        "Implemented hybrid",
    ),
    (
        "User Simulator",
        "Dung simulator/offline interaction de train/evaluate policy.",
        "MovieLens-derived user profiles, positive/negative items va attribute feedback.",
        "Implemented demo-scale",
    ),
    (
        "Evaluation",
        "SR@T, average turns, ablation tren dataset paper.",
        "Offline evaluation + negative sampler ablation; live session panel cho demo.",
        "Implemented, small scale",
    ),
    (
        "Scale/Reproduction",
        "Full experimental protocol tren KB4REC/Freebase.",
        "Demo report: TMDB cap 500 movies, bounded candidate pool, small user sample.",
        "Limitation",
    ),
]


def set_text(shape, text, font_size=18, bold=False, color=BLACK, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.55))
    set_text(title_box, title, 25, True, NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.90), Inches(12.0), Inches(0.35))
        set_text(sub, subtitle, 11, False, GRAY)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.22), Inches(12.25), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(11.6), Inches(7.0), Inches(1.1), Inches(0.25))
    set_text(box, f"{idx}", 9, False, GRAY, PP_ALIGN.RIGHT)


def add_note(slide, text, x, y, w, h, color=LIGHT, border=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK


def add_bullets(slide, items, x, y, w, h, font_size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(5)


def add_table(slide, rows, col_widths, x, y, row_h=0.42, font_size=9.5):
    table_shape = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(x),
        Inches(y),
        Inches(sum(col_widths)),
        Inches(row_h * len(rows)),
    )
    table = table_shape.table
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(row_h if r else row_h * 0.85)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif c == len(row) - 1 and "Limitation" in value:
                cell.fill.fore_color.rgb = RGBColor(255, 243, 224)
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(250, 251, 253)
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = value
            p.font.name = "Arial"
            p.font.size = Pt(font_size if r else font_size + 0.5)
            p.font.bold = r == 0
            p.font.color.rgb = WHITE if r == 0 else BLACK
            p.alignment = PP_ALIGN.LEFT


def add_metric_card(slide, x, y, w, h, label, value, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.25)
    v = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(0.35))
    set_text(v, value, 20, True, color)
    l = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.55), Inches(w - 0.24), Inches(0.25))
    set_text(l, label, 9, False, GRAY)


def slide_comparison_1(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Paper vs Code: Component Mapping", "Nhung thanh phan paper da duoc hien thuc trong project demo")
    rows = [("Paper component", "KGenSam paper", "Code demo", "Status")]
    rows += COMPARISON_ROWS[:5]
    add_table(slide, rows, [1.7, 3.45, 4.25, 1.65], 0.55, 1.48, row_h=0.82, font_size=8.6)
    add_footer(slide, idx)


def slide_comparison_2(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Paper vs Code: Training, Evaluation, Scope", "Phan nay nen dua vao bao cao de noi ro muc do reproduce paper")
    rows = [("Paper component", "KGenSam paper", "Code demo", "Status")]
    rows += COMPARISON_ROWS[5:]
    add_table(slide, rows, [1.7, 3.45, 4.25, 1.65], 0.55, 1.48, row_h=0.82, font_size=8.6)
    add_note(
        slide,
        "Ket luan de noi khi thuyet trinh:\nProject khong claim reproduce 100% paper. Project implement core CRS loop theo tinh than KGenSam: KG + Ask/Recommend policy + Active Sampler + Negative Sampler + evaluation demo.",
        0.78,
        6.0,
        11.35,
        0.85,
        color=RGBColor(240, 248, 255),
        border=TEAL,
    )
    add_footer(slide, idx)


def slide_workflow(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Workflow Demo Theo KGenSam", "Moi turn: policy quyet dinh ASK hoac RECOMMEND, user tra feedback nhi phan")
    boxes = [
        ("User start", 0.85, 2.0, BLUE),
        ("Interact Policy\nrollout_dqn", 3.0, 2.0, TEAL),
        ("ASK\nActive Sampler", 5.35, 1.35, ORANGE),
        ("Binary feedback\nAccept / Reject attribute", 7.75, 1.35, BLUE),
        ("RECOMMEND\nHybrid ranker", 5.35, 3.4, ORANGE),
        ("Item feedback\nLooks good / Not for me", 7.75, 3.4, BLUE),
        ("Stop success\nor max T", 10.35, 2.45, TEAL),
    ]
    for text, x, y, color in boxes:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.75), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(247, 250, 252)
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        set_text(shape, text, 12, True if "\n" not in text else False, BLACK, PP_ALIGN.CENTER)
    arrows = [
        (2.55, 2.38, 0.42, 0),
        (4.75, 2.15, 0.62, -0.55),
        (4.75, 2.68, 0.62, 0.75),
        (7.1, 1.75, 0.6, 0),
        (7.1, 3.8, 0.6, 0),
        (9.5, 1.75, 0.85, 0.75),
        (9.5, 3.8, 0.85, -0.65),
    ]
    for x, y, w, h in arrows:
        line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
        line.line.color.rgb = GRAY
        line.line.width = Pt(1.5)
    add_note(
        slide,
        "Demo setting: min_ask_turns = 3, max_turns = 5.\nLy do: hien ro workflow Ask/Feedback truoc khi recommend, tranh recommend qua som trong luc thuyet trinh.",
        0.85,
        5.35,
        11.15,
        0.75,
        color=RGBColor(255, 250, 240),
        border=ORANGE,
    )
    add_footer(slide, idx)


def slide_architecture(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Code Architecture Da Implement", "Mapping tu workflow paper sang cac module trong backend/frontend")
    rows = [
        ("Layer", "Module/File", "Vai tro trong demo"),
        ("Frontend", "src/ui/chat.js", "Chat flow, Ask/Recommend UI, user feedback buttons, live session event"),
        ("Visualization", "src/ui/graphViz.js", "Ve KG ben phai: user, attributes, recommendations, edges ask/recommend"),
        ("API", "backend/app/main.py", "Conversation endpoints, evaluation, report summary, TMDB cache handling"),
        ("Policy", "backend/app/interact_policy.py", "rollout_dqn quyet dinh ASK/RECOMMEND"),
        ("Active Sampler", "backend/app/active_sampler.py", "Chon attribute question bang entropy/centrality/GCN policy"),
        ("Recommender", "backend/app/recommender.py", "FM/BPR + KG propagation + embedding similarity"),
        ("Negative Sampler", "backend/app/negative_sampler.py", "learned_linear_policy negative sampler"),
    ]
    add_table(slide, rows, [1.55, 3.15, 7.05], 0.62, 1.45, row_h=0.58, font_size=9)
    add_footer(slide, idx)


def slide_dataset(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Dataset va Knowledge Graph Trong Demo", "Thay Freebase bang MovieLens + TMDB external KG")
    add_metric_card(slide, 0.75, 1.45, 2.3, 0.9, "MovieLens movies", "9,742", BLUE)
    add_metric_card(slide, 3.35, 1.45, 2.3, 0.9, "MovieLens triples", "31,779", TEAL)
    add_metric_card(slide, 5.95, 1.45, 2.3, 0.9, "TMDB enriched movies", "500", ORANGE)
    add_metric_card(slide, 8.55, 1.45, 2.3, 0.9, "TMDB triples", "5,729", BLUE)
    add_metric_card(slide, 3.0, 2.65, 2.7, 0.9, "Final KG entities", "12,797", NAVY)
    add_metric_card(slide, 6.05, 2.65, 2.7, 0.9, "Final KG triples", "37,508", NAVY)
    add_bullets(
        slide,
        [
            "MovieLens ratings.csv: user-item interaction.",
            "MovieLens movies.csv: title, genre, release year.",
            "MovieLens links.csv: mapping movieId -> tmdbId.",
            "TMDB API/cache: director, cast, keywords, language.",
            "Local curated KG da bo; demo chi giu standard dataset + external KG.",
        ],
        1.0,
        4.05,
        10.8,
        1.45,
        font_size=15,
    )
    add_footer(slide, idx)


def slide_evaluation(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Evaluation va Ablation Trong Code Demo", "Metric dung de chung minh pipeline chay duoc, khong claim full reproduction")
    add_metric_card(slide, 0.9, 1.55, 2.35, 0.95, "SR@T", "0.60", BLUE)
    add_metric_card(slide, 3.55, 1.55, 2.35, 0.95, "Average turns", "1.00", TEAL)
    add_metric_card(slide, 6.2, 1.55, 2.35, 0.95, "Average asks", "1.00", ORANGE)
    add_metric_card(slide, 8.85, 1.55, 2.35, 0.95, "Average recommends", "1.00", NAVY)
    add_note(
        slide,
        "Offline Evaluation config:\nmax_users = 5, max_turns = 5, candidate_pool = 150, seed = 42",
        0.9,
        3.0,
        5.35,
        1.0,
        border=BLUE,
    )
    add_note(
        slide,
        "Ablation:\nlearned_negative_current vs random_negative_baseline\nMuc dich: cho thay sampling strategy co the thay doi va do anh huong len CRS pipeline.",
        6.65,
        3.0,
        5.4,
        1.0,
        border=ORANGE,
    )
    add_bullets(
        slide,
        [
            "Run: chay offline simulator va hien SR@T/turns/asks/recommends.",
            "Ablation: so sanh negative sampler hien tai voi random baseline.",
            "Live Session: theo doi lua chon that cua user tren browser, khac voi offline simulator.",
        ],
        1.0,
        4.65,
        10.7,
        1.2,
        font_size=14,
    )
    add_footer(slide, idx)


def slide_limitations(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Han Che So Voi Paper va Cach Noi Trong Bao Cao", "Day la slide nen de cuoi truoc Q&A")
    rows = [
        ("Han che", "Giai thich ngan de dua vao slide/report"),
        ("Khong reproduce 100% paper", "Project la demo implementation, khong phai full experimental reproduction tren KB4REC."),
        ("External KG khac paper", "Paper dung Freebase; demo dung TMDB vi Freebase da discontinued va TMDB phu hop movie domain."),
        ("RL training demo-scale", "Interact/Active policies train bang MovieLens-derived offline simulator rollouts, chua phai full online RL."),
        ("TMDB coverage gioi han", "TMDB_MAX_MOVIES = 500 de startup va demo on dinh."),
        ("Evaluation nho", "Dung bounded candidate pool va it users; metric chi mang tinh tham chieu demo."),
        ("Ablation chua thong ke", "Chua chay nhieu seed/dataset/statistical test nen khong ket luan manh nhu paper."),
    ]
    add_table(slide, rows, [2.55, 8.8], 0.72, 1.45, row_h=0.63, font_size=9.5)
    add_note(
        slide,
        "Cau chot nen noi:\nWe implement the main KGenSam-inspired CRS loop and demonstrate the core mechanism, while leaving full-scale reproduction and large benchmark evaluation as future work.",
        0.92,
        6.25,
        10.95,
        0.65,
        color=RGBColor(240, 248, 255),
        border=TEAL,
    )
    add_footer(slide, idx)


def write_markdown():
    lines = [
        "# Bang so sanh Paper KGenSam vs Code Demo",
        "",
        "| Component | Paper KGenSam | Code demo | Muc do |",
        "|---|---|---|---|",
    ]
    for row in COMPARISON_ROWS:
        lines.append(f"| {row[0]} | {row[1]} | {row[2]} | {row[3]} |")
    lines += [
        "",
        "## Cau chot dua vao slide/report",
        "",
        "Project khong claim reproduce 100% paper. Project implement core CRS loop theo tinh than KGenSam: KG + Interact Policy ASK/RECOMMEND + Active Sampler + Negative Sampler + hybrid recommender + binary feedback + offline evaluation.",
    ]
    MD_OUT.write_text("\n".join(lines), encoding="utf-8")


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    prs = Presentation(str(SRC))
    start_idx = len(prs.slides) + 1
    slide_comparison_1(prs, start_idx)
    slide_comparison_2(prs, start_idx + 1)
    slide_workflow(prs, start_idx + 2)
    slide_architecture(prs, start_idx + 3)
    slide_dataset(prs, start_idx + 4)
    slide_evaluation(prs, start_idx + 5)
    slide_limitations(prs, start_idx + 6)
    prs.save(str(OUT))
    write_markdown()
    print(f"saved {OUT}")
    print(f"saved {MD_OUT}")
    print(f"slides {len(prs.slides)}")


if __name__ == "__main__":
    main()
